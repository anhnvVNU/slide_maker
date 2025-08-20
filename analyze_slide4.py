#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu

def analyze_slide_formatting(pptx_path, slide_index=3):
    """Analyze formatting of a specific slide (0-based index)"""
    print(f"Analyzing slide {slide_index + 1} from: {pptx_path}")
    
    # Load presentation
    prs = Presentation(str(pptx_path))
    
    if slide_index >= len(prs.slides):
        print(f"Error: Slide {slide_index + 1} not found. Presentation has {len(prs.slides)} slides.")
        return
    
    slide = prs.slides[slide_index]
    
    analysis = {
        "slide_number": slide_index + 1,
        "slide_layout": slide.slide_layout.name,
        "shapes_count": len(slide.shapes),
        "shapes": []
    }
    
    print(f"\nSlide {slide_index + 1} Layout: {slide.slide_layout.name}")
    print(f"Number of shapes: {len(slide.shapes)}")
    print("=" * 60)
    
    for i, shape in enumerate(slide.shapes):
        shape_info = {
            "shape_index": i,
            "shape_type": str(shape.shape_type),
            "has_text_frame": hasattr(shape, 'text_frame'),
            "position": {
                "x": int(shape.left) if hasattr(shape, 'left') else None,
                "y": int(shape.top) if hasattr(shape, 'top') else None
            },
            "size": {
                "width": int(shape.width) if hasattr(shape, 'width') else None,
                "height": int(shape.height) if hasattr(shape, 'height') else None
            }
        }
        
        print(f"\nShape {i}:")
        print(f"  Type: {shape.shape_type}")
        print(f"  Position: x={shape_info['position']['x']}, y={shape_info['position']['y']}")
        print(f"  Size: {shape_info['size']['width']} x {shape_info['size']['height']}")
        
        # Analyze text frame if exists
        if hasattr(shape, 'text_frame') and shape.has_text_frame:
            text_frame = shape.text_frame
            shape_info["text_frame"] = {
                "text": shape.text,
                "paragraphs_count": len(text_frame.paragraphs),
                "paragraphs": []
            }
            
            print(f"  Text: '{shape.text}'")
            print(f"  Paragraphs: {len(text_frame.paragraphs)}")
            
            for j, paragraph in enumerate(text_frame.paragraphs):
                para_info = {
                    "paragraph_index": j,
                    "text": paragraph.text,
                    "level": paragraph.level,
                    "alignment": str(paragraph.alignment),
                    "runs_count": len(paragraph.runs)
                }
                
                # Analyze spacing
                try:
                    para_info["space_before"] = {
                        "value": float(paragraph.space_before.pt) if paragraph.space_before else None,
                        "unit": "pt"
                    }
                except:
                    para_info["space_before"] = None
                
                try:
                    para_info["space_after"] = {
                        "value": float(paragraph.space_after.pt) if paragraph.space_after else None,
                        "unit": "pt"
                    }
                except:
                    para_info["space_after"] = None
                
                try:
                    para_info["line_spacing"] = {
                        "value": float(paragraph.line_spacing) if paragraph.line_spacing else None
                    }
                except:
                    para_info["line_spacing"] = None
                
                # Analyze bullet formatting via XML
                try:
                    from lxml import etree
                    p_elem = paragraph._p
                    pPr = p_elem.find('.//a:pPr', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    
                    if pPr is not None:
                        para_info["xml_properties"] = {}
                        
                        # Get indent and margin
                        indent = pPr.get('indent')
                        marL = pPr.get('marL') 
                        
                        if indent:
                            para_info["xml_properties"]["indent"] = {
                                "value": int(indent),
                                "inches": int(indent) / 914400,
                                "unit": "EMU"
                            }
                        
                        if marL:
                            para_info["xml_properties"]["marL"] = {
                                "value": int(marL),
                                "inches": int(marL) / 914400,
                                "unit": "EMU"
                            }
                        
                        # Check for bullet character
                        buChar = pPr.find('.//a:buChar', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                        if buChar is not None:
                            para_info["xml_properties"]["bullet_char"] = buChar.get('char')
                        
                        # Check for numbering
                        buAutoNum = pPr.find('.//a:buAutoNum', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                        if buAutoNum is not None:
                            para_info["xml_properties"]["numbering"] = {
                                "type": buAutoNum.get('type'),
                                "startAt": buAutoNum.get('startAt')
                            }
                
                except Exception as e:
                    para_info["xml_error"] = str(e)
                
                print(f"    Para {j}: '{paragraph.text}' (level {paragraph.level})")
                if para_info.get("space_before"):
                    print(f"      Space before: {para_info['space_before']['value']}pt")
                if para_info.get("space_after"):
                    print(f"      Space after: {para_info['space_after']['value']}pt")
                if para_info.get("line_spacing"):
                    print(f"      Line spacing: {para_info['line_spacing']['value']}")
                if para_info.get("xml_properties"):
                    xml_props = para_info["xml_properties"]
                    if "indent" in xml_props:
                        print(f"      Indent: {xml_props['indent']['value']} EMU ({xml_props['indent']['inches']:.3f}\")")
                    if "marL" in xml_props:
                        print(f"      Left margin: {xml_props['marL']['value']} EMU ({xml_props['marL']['inches']:.3f}\")")
                    if "bullet_char" in xml_props:
                        print(f"      Bullet char: '{xml_props['bullet_char']}'")
                    if "numbering" in xml_props:
                        print(f"      Numbering: {xml_props['numbering']}")
                
                # Analyze runs (font formatting)
                para_info["runs"] = []
                for k, run in enumerate(paragraph.runs):
                    run_info = {
                        "run_index": k,
                        "text": run.text,
                        "font": {}
                    }
                    
                    font = run.font
                    if font.name:
                        run_info["font"]["name"] = font.name
                    try:
                        if font.size:
                            run_info["font"]["size_pt"] = font.size.pt
                    except:
                        pass
                    if font.bold is not None:
                        run_info["font"]["bold"] = font.bold
                    if font.italic is not None:
                        run_info["font"]["italic"] = font.italic
                    try:
                        if font.color and hasattr(font.color, 'rgb') and font.color.rgb:
                            rgb = font.color.rgb
                            run_info["font"]["color"] = [rgb.red, rgb.green, rgb.blue]
                    except Exception as e:
                        run_info["font"]["color_error"] = str(e)
                    
                    para_info["runs"].append(run_info)
                    print(f"      Run {k}: '{run.text}' - Font: {run_info['font']}")
                
                shape_info["text_frame"]["paragraphs"].append(para_info)
        
        analysis["shapes"].append(shape_info)
        print("-" * 40)
    
    return analysis

def compare_slides(file1_path, file2_path, slide_index=3):
    """Compare slide 4 between two PowerPoint files"""
    print(f"COMPARING SLIDE {slide_index + 1}:")
    print(f"File 1 (Template): {file1_path}")
    print(f"File 2 (Output):   {file2_path}")
    print("=" * 80)
    
    analysis1 = analyze_slide_formatting(file1_path, slide_index)
    print("\n" + "="*80)
    analysis2 = analyze_slide_formatting(file2_path, slide_index)
    
    if not analysis1 or not analysis2:
        print("Error: Could not analyze one or both files")
        return
    
    # Save both analyses
    with open("template_slide4_analysis.json", 'w', encoding='utf-8') as f:
        json.dump(analysis1, f, ensure_ascii=False, indent=2)
    
    with open("output_slide4_analysis.json", 'w', encoding='utf-8') as f:
        json.dump(analysis2, f, ensure_ascii=False, indent=2)
    
    print(f"\n{'='*80}")
    print("COMPARISON SUMMARY:")
    print(f"{'='*80}")
    
    # Compare main content shape (position x=626250)
    template_content = None
    output_content = None
    
    for shape in analysis1["shapes"]:
        if shape.get("position", {}).get("x") == 626250:
            template_content = shape
            break
    
    for shape in analysis2["shapes"]:
        if shape.get("position", {}).get("x") == 626250:
            output_content = shape
            break
    
    if template_content and output_content:
        print(f"\nMAIN CONTENT SHAPE COMPARISON (x=626250):")
        print("-" * 50)
        
        t_paras = template_content.get("text_frame", {}).get("paragraphs", [])
        o_paras = output_content.get("text_frame", {}).get("paragraphs", [])
        
        print(f"Template paragraphs: {len(t_paras)}")
        print(f"Output paragraphs:   {len(o_paras)}")
        
        # Compare first few paragraphs
        max_compare = min(len(t_paras), len(o_paras), 5)
        
        for i in range(max_compare):
            t_para = t_paras[i]
            o_para = o_paras[i]
            
            print(f"\nParagraph {i}:")
            print(f"  Template: '{t_para['text'][:40]}...' (level {t_para['level']})")
            print(f"  Output:   '{o_para['text'][:40]}...' (level {o_para['level']})")
            
            # Compare spacing
            t_before = t_para.get("space_before", {}).get("value")
            o_before = o_para.get("space_before", {}).get("value")
            t_after = t_para.get("space_after", {}).get("value")
            o_after = o_para.get("space_after", {}).get("value")
            
            print(f"  Space before: Template={t_before}pt, Output={o_before}pt")
            print(f"  Space after:  Template={t_after}pt, Output={o_after}pt")
            
            # Compare XML properties
            t_xml = t_para.get("xml_properties", {})
            o_xml = o_para.get("xml_properties", {})
            
            if t_xml or o_xml:
                print(f"  XML Properties:")
                
                # Compare bullets
                t_bullet = t_xml.get("bullet_char")
                o_bullet = o_xml.get("bullet_char")
                if t_bullet or o_bullet:
                    match = "✅" if t_bullet == o_bullet else "❌"
                    print(f"    Bullet char: Template='{t_bullet}', Output='{o_bullet}' {match}")
                
                # Compare indent
                t_indent = t_xml.get("indent", {}).get("inches")
                o_indent = o_xml.get("indent", {}).get("inches")
                if t_indent is not None or o_indent is not None:
                    match = "✅" if abs((t_indent or 0) - (o_indent or 0)) < 0.001 else "❌"
                    print(f"    Indent: Template={t_indent:.3f}\", Output={o_indent:.3f}\" {match}")
                
                # Compare margin
                t_marL = t_xml.get("marL", {}).get("inches")
                o_marL = o_xml.get("marL", {}).get("inches")
                if t_marL is not None or o_marL is not None:
                    match = "✅" if abs((t_marL or 0) - (o_marL or 0)) < 0.001 else "❌"
                    print(f"    MarginL: Template={t_marL:.3f}\", Output={o_marL:.3f}\" {match}")

def main():
    template_path = Path("/media/vunv/DATA/slide_maker_4/slide_maker_2/VuNV/slide_maker/template/20250813.pptx")
    output_path = Path("/media/vunv/DATA/slide_maker_4/slide_maker_2/VuNV/slide_maker/output/final_presentation.pptx")
    
    # Check if files exist
    if not template_path.exists():
        print(f"Error: Template file not found: {template_path}")
        return
    
    if not output_path.exists():
        print(f"Error: Output file not found: {output_path}")
        return
    
    # Compare slide 4 of both files
    compare_slides(template_path, output_path, slide_index=3)

if __name__ == "__main__":
    main()
