#!/usr/bin/env python3
"""
Main script for Slide Maker
Complete pipeline: Raw text → Slide Script → PowerPoint presentation
"""

import sys
import json
import io
from pathlib import Path

# Add code directory to path for imports
sys.path.append(str(Path(__file__).parent / "code"))

from script_generator import SlideScriptGenerator as ScriptGenerator
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def add_text_to_slide(slide, shape_data):
    """Add a text shape to the slide based on shape data from JSON"""
    position = shape_data.get('position', {})
    size = shape_data.get('size', {})
    
    # Convert EMU values
    left = Emu(position.get('x', 0))
    top = Emu(position.get('y', 0))
    width = Emu(size.get('width', 1000000))
    height = Emu(size.get('height', 500000))
    
    # Add textbox
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    
    # Configure text frame
    text_frame.word_wrap = True
    text_frame.margin_left = Emu(45720)
    text_frame.margin_right = Emu(45720)
    text_frame.margin_top = Emu(22860)
    text_frame.margin_bottom = Emu(22860)
    
    # Add paragraphs
    paragraphs_data = shape_data.get('paragraphs', [])
    
    for i, para_data in enumerate(paragraphs_data):
        # Add paragraph
        if i == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        
        # Set alignment
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        alignment = para_data.get('alignment', 'left')
        if alignment in alignment_map:
            paragraph.alignment = alignment_map[alignment]
        
        # Set line spacing and paragraph spacing
        line_spacing = para_data.get('line_spacing', {})
        if line_spacing:
            spacing_value = line_spacing.get('value', 1.0)
            spacing_type = line_spacing.get('type', 'multiple')
            
            if spacing_type == 'points':
                paragraph.space_after = Pt(spacing_value * 12)
                if spacing_value > 1.0:
                    paragraph.line_spacing = spacing_value
            else:
                paragraph.line_spacing = spacing_value
        
        # Set space before if specified
        space_before = para_data.get('space_before_pt', para_data.get('space_before', 0))
        if space_before:
            paragraph.space_before = Pt(space_before)
        
        # Set space after if specified  
        space_after = para_data.get('space_after_pt', para_data.get('space_after', 0))
        if space_after:
            paragraph.space_after = Pt(space_after)
        elif i < len(paragraphs_data) - 1:
            # Add default spacing after each paragraph (except last) if not explicitly set
            paragraph.space_after = Pt(12)
        
        # Set paragraph level for bullets
        level = para_data.get('level', 0)
        paragraph.level = level
        
        # Handle bullets
        bullet_data = para_data.get('bullet', {})
        if bullet_data:
            bullet_type = bullet_data.get('type', 'none')
            if bullet_type == 'bullet':
                bullet_char = bullet_data.get('char', '•')
                if bullet_char and len(bullet_char) > 0:
                    # Use level from JSON data, don't override based on character
                    paragraph.level = level
                    
                    # Enable bullet formatting by accessing XML
                    try:
                        from lxml import etree
                        # Access the paragraph XML element
                        p_elem = paragraph._p
                        
                        # Find or create pPr (paragraph properties)
                        pPr = p_elem.find('.//a:pPr', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                        if pPr is None:
                            pPr = etree.SubElement(p_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                        
                        # Set indentation based on level (matching template)
                        if level == 0:
                            pPr.set('indent', '-317500')  # Negative indent pulls bullet left
                            pPr.set('marL', '457200')     # Left margin for text content
                        elif level == 1:
                            pPr.set('indent', '-317500')  # Same negative indent for both levels
                            pPr.set('marL', '914400')     # Larger left margin for level 1
                        
                        # Add bullet character formatting
                        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                        buChar.set('char', bullet_char)
                        
                    except Exception as e:
                        # If XML manipulation fails, fall back to just setting level
                        print(f"Warning: Could not set bullet format: {e}")
                        pass
            elif bullet_type == 'numbered':
                # Handle numbered lists
                # Set the paragraph as a numbered list
                paragraph.level = level  # Use the level from para_data
                
                # Enable numbering by accessing the XML element
                try:
                    from lxml import etree
                    # Access the paragraph XML element
                    p_elem = paragraph._p
                    
                    # Find or create pPr (paragraph properties)
                    pPr = p_elem.find('.//a:pPr', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    if pPr is None:
                        pPr = etree.SubElement(p_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                    
                    # Add numbering format
                    buAutoNum = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                    buAutoNum.set('type', 'arabicPeriod')  # Creates 1. 2. 3. format
                    buAutoNum.set('startAt', '1')
                    
                except Exception as e:
                    # If XML manipulation fails, fall back to just setting level
                    print(f"Warning: Could not set numbering format: {e}")
                    pass
        
        # Add text runs
        runs_data = para_data.get('runs', [])
        for j, run_data in enumerate(runs_data):
            # If there are multiple runs and we want them on separate lines,
            # add a line break between them
            if j > 0:
                # Add a line break before this run (except for the first run)
                run = paragraph.add_run()
                run.text = '\n'
            
            run = paragraph.add_run()
            run.text = run_data.get('text', '')
            
            # Format font
            font_data = run_data.get('font', {})
            font = run.font
            
            # Set font name
            font_name = font_data.get('name', 'Arial')
            if font_name:
                font.name = font_name
            
            # Set font size
            font_size = font_data.get('size_pt')
            if font_size:
                font.size = Pt(font_size)
            else:
                font.size = Pt(12)
            
            # Set bold
            if 'bold' in font_data:
                font.bold = font_data['bold']
            
            # Set italic
            if 'italic' in font_data:
                font.italic = font_data['italic']
            
            # Set color
            color = font_data.get('color')
            if color and isinstance(color, list) and len(color) >= 3:
                font.color.rgb = RGBColor(color[0], color[1], color[2])


def main():
    """
    Main function to run the complete slide generation pipeline
    """
    print("="*70)
    print(" SLIDE MAKER - Presentation Generation Pipeline")
    print("="*70)
    
    # Define paths
    base_dir = Path(__file__).parent
    input_file = base_dir / "data" / "raw_data.txt"
    output_dir = base_dir / "output"
    script_file = output_dir / "slide_script.json"
    template_path = base_dir / "template" / "template.pptx"
    
    # Create output directory
    output_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        # STEP 1: Generate slide script from raw data
        print("\n📝 STEP 1: Generating Slide Script")
        print("-" * 50)
        
        generator = ScriptGenerator()
        
        # Read input file
        print(f"Reading input from: {input_file}")
        with open(input_file, 'r', encoding='utf-8') as f:
            raw_text = f.read()
        
        # Save raw text to represent.txt for the generator
        represent_file = output_dir / "represent.txt"
        with open(represent_file, 'w', encoding='utf-8') as f:
            f.write(raw_text)
        print(f"✓ Saved raw text to: {represent_file.name}")
        
        # Generate slides using the generator's methods
        print("Generating slide structure...")
        slides = generator.generate_presentation_script()
        print(f"✓ Generated {len(slides)} slides")
        
        # Save script
        script_path = generator.save_script(slides)
        print(f"✓ Saved slide script to: {script_file.name}")
        
        # STEP 2: Build presentation from script and template
        print("\n🎨 STEP 2: Building PowerPoint Presentation")
        print("-" * 50)
        
        # Load template
        print(f"Loading master template: {template_path.name}")
        prs = Presentation(str(template_path))
        
        # Map layouts from template slides
        layout_map = {}
        if len(prs.slides) >= 4:
            layout_map['cover'] = prs.slides[0].slide_layout
            layout_map['agenda_slide'] = prs.slides[1].slide_layout
            layout_map['subtitle'] = prs.slides[2].slide_layout
            layout_map['section_divider'] = prs.slides[2].slide_layout
            layout_map['content'] = prs.slides[3].slide_layout
            print("✓ Identified layouts from template")
        
        original_slide_count = len(prs.slides)
        
        # Load generated script
        with open(script_file, 'r', encoding='utf-8') as f:
            slides_data = json.load(f)
        
        # Create slides with content
        print(f"Creating {len(slides_data)} slides...")
        for i, slide_data in enumerate(slides_data, 1):
            slide_type = slide_data.get('slide_type', 'content')
            
            # Get appropriate layout
            layout = layout_map.get(slide_type, layout_map.get('content'))
            
            # Add new slide
            new_slide = prs.slides.add_slide(layout)
            
            # Map template slide index
            template_slide_idx = None
            if slide_type == 'cover':
                template_slide_idx = 0
            elif slide_type == 'agenda_slide':
                template_slide_idx = 1
            elif slide_type in ['subtitle', 'section_divider']:
                template_slide_idx = 2
            elif slide_type == 'content':
                template_slide_idx = 3
            
            # Copy images from template
            if template_slide_idx is not None and template_slide_idx < original_slide_count:
                template_slide = prs.slides[template_slide_idx]
                
                for shape in template_slide.shapes:
                    # Copy images
                    if shape.shape_type == 13:  # Picture
                        try:
                            image_bytes = shape.image.blob
                            new_slide.shapes.add_picture(
                                io.BytesIO(image_bytes), 
                                shape.left, shape.top, 
                                shape.width, shape.height
                            )
                        except:
                            pass
            
            # Remove empty placeholders
            shapes_to_remove = []
            for shape in new_slide.shapes:
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    try:
                        if shape.placeholder_format and (not shape.text or not shape.text.strip()):
                            shapes_to_remove.append(shape)
                    except:
                        if not shape.text or not shape.text.strip():
                            shapes_to_remove.append(shape)
            
            for shape in shapes_to_remove:
                sp = shape._element
                sp.getparent().remove(sp)
            
            # Add content from JSON
            shapes_data = slide_data.get('shapes', [])
            for shape_data in shapes_data:
                if shape_data.get('type') == 'text':
                    add_text_to_slide(new_slide, shape_data)
            
            print(f"  ✓ Slide {i:2d}: {slide_type}")
        
        # Remove original template slides
        print(f"Removing {original_slide_count} template slides...")
        for i in range(original_slide_count):
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            if slides:
                xml_slides.remove(slides[0])
        
        # Save final presentation
        output_file = output_dir / "final_presentation.pptx"
        prs.save(str(output_file))
        print(f"✓ Saved presentation: {output_file.name}")
        
        # Success summary
        print("\n" + "="*70)
        print(" ✅ SUCCESS! Presentation Generated")
        print("="*70)
        print(f" 📄 Input:        {input_file.name}")
        print(f" 📋 Script:       {script_file.name}")
        print(f" 🎯 Presentation: {output_file.name}")
        print(f" 📊 Total slides: {len(slides_data)}")
        print("="*70)
        
    except Exception as e:
        print(f"\n❌ Error: {e}")
        print("\nTroubleshooting:")
        print("1. Check input file exists: data/raw_data.txt")
        print("2. Ensure template exists: template/template.pptx")
        print("3. Verify all dependencies are installed")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
