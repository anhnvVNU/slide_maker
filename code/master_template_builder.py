#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import json
import io
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_text_to_slide(slide, shape_data):
    """Add a text shape to the slide based on shape data from JSON"""
    position = shape_data.get('position', {})
    size = shape_data.get('size', {})
    
    # Convert EMU values
    left = Emu(position.get('x', 0))
    top = Emu(position.get('y', 0))
    width = Emu(size.get('width', 1000000))
    height = Emu(size.get('height', 500000))
    
    # Add textbox
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    
    # Handle fill color if specified
    fill_data = shape_data.get('fill', {})
    if fill_data:
        fill_type = fill_data.get('type')
        if fill_type == 'solid':
            color = fill_data.get('color')
            if color and isinstance(color, list) and len(color) >= 3:
                textbox.fill.solid()
                textbox.fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])
    
    # Configure text frame
    text_frame.word_wrap = True
    text_frame.margin_left = Emu(45720)
    text_frame.margin_right = Emu(45720)
    text_frame.margin_top = Emu(22860)
    text_frame.margin_bottom = Emu(22860)
    
    # Add paragraphs
    paragraphs_data = shape_data.get('paragraphs', [])
    
    for i, para_data in enumerate(paragraphs_data):
        # Add paragraph
        if i == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            # Add a new paragraph (this automatically adds a line break)
            paragraph = text_frame.add_paragraph()
        
        # Set alignment
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        alignment = para_data.get('alignment', 'left')
        if alignment in alignment_map:
            paragraph.alignment = alignment_map[alignment]
        
        # Set line spacing - KHÔNG tạo space_after để match template gốc
        line_spacing = para_data.get('line_spacing', {})
        if line_spacing:
            spacing_value = line_spacing.get('value', 1.0)
            spacing_type = line_spacing.get('type', 'multiple')
             
            # Template gốc chỉ có line_spacing, KHÔNG có space_after
            paragraph.line_spacing = spacing_value
        
        # Set space before if specified
        space_before = para_data.get('space_before_pt', para_data.get('space_before', 0))
        if space_before:
            paragraph.space_before = Pt(space_before)
        
        # Set space after if specified  
        space_after = para_data.get('space_after_pt', para_data.get('space_after', 0))
        if space_after:
            paragraph.space_after = Pt(space_after)
        # elif i < len(paragraphs_data) - 1:
        #     # Add default spacing after each paragraph (except last) if not explicitly set
        #     # Template gốc KHÔNG CÓ spacing - tắt để match với gốc
        #     paragraph.space_after = Pt(6)
        
        # Set paragraph level for bullets
        level = para_data.get('level', 0)
        paragraph.level = level
        
        # Handle bullets
        bullet_data = para_data.get('bullet', {})
        if bullet_data:
            bullet_type = bullet_data.get('type', 'none')
            if bullet_type == 'bullet':
                bullet_char = bullet_data.get('char', '•')
                if bullet_char and len(bullet_char) > 0:
                    # Use level from JSON data
                    paragraph.level = level
                    
                    # Enable bullet formatting by accessing XML
                    try:
                        from lxml import etree
                        # Access the paragraph XML element
                        p_elem = paragraph._p
                        
                        # Find or create pPr (paragraph properties)
                        pPr = p_elem.find('.//a:pPr', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                        if pPr is None:
                            pPr = etree.SubElement(p_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                        
                        # Set indentation based on level (matching template)
                        if level == 0:
                            pPr.set('indent', '-317500')  # Negative indent pulls bullet left
                            pPr.set('marL', '457200')     # Left margin for text content
                        elif level == 1:
                            pPr.set('indent', '-317500')  # Same negative indent for both levels
                            pPr.set('marL', '914400')     # Larger left margin for level 1
                        
                        # Add bullet character formatting
                        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                        buChar.set('char', bullet_char)
                        
                    except Exception as e:
                        # If XML manipulation fails, fall back to just setting level
                        print(f"Warning: Could not set bullet format: {e}")
                        pass
            elif bullet_type == 'numbered':
                # Handle numbered lists
                # Set the paragraph as a numbered list
                paragraph.level = level  # Use the level from para_data
                
                # Enable numbering by accessing the XML element
                try:
                    from lxml import etree
                    # Access the paragraph XML element
                    p_elem = paragraph._p
                    
                    # Find or create pPr (paragraph properties)
                    pPr = p_elem.find('.//a:pPr', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    if pPr is None:
                        pPr = etree.SubElement(p_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                    
                    # Set indentation for numbered lists (matching template)
                    # Template gốc: indent=-336550 EMU (-26.5pt), marL=457200 EMU (36.0pt)
                    pPr.set('indent', '-336550')  # Negative indent pulls number left
                    pPr.set('marL', '457200')     # Left margin for text content (0.5 inch)
                    
                    # Add numbering format
                    buAutoNum = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                    buAutoNum.set('type', 'arabicPeriod')  # Creates 1. 2. 3. format
                    buAutoNum.set('startAt', '1')
                    
                except Exception as e:
                    # If XML manipulation fails, fall back to just setting level
                    print(f"Warning: Could not set numbering format: {e}")
                    pass
        
        # Add text runs
        runs_data = para_data.get('runs', [])
        for j, run_data in enumerate(runs_data):
            # If there are multiple runs and we want them on separate lines,
            # add a line break between them
            if j > 0:
                # Add a line break before this run (except for the first run)
                run = paragraph.add_run()
                run.text = '\n'
            
            run = paragraph.add_run()
            run.text = run_data.get('text', '')
            
            # Format font
            font_data = run_data.get('font', {})
            font = run.font
            
            # Set font name
            font_name = font_data.get('name', 'Arial')
            if font_name:
                font.name = font_name
            
            # Set font size
            font_size = font_data.get('size_pt')
            if font_size:
                font.size = Pt(font_size)
            else:
                font.size = Pt(12)
            
            # Set bold
            if 'bold' in font_data:
                font.bold = font_data['bold']
            
            # Set italic
            if 'italic' in font_data:
                font.italic = font_data['italic']
            
            # Set color
            color = font_data.get('color')
            if color and isinstance(color, list) and len(color) >= 3:
                font.color.rgb = RGBColor(color[0], color[1], color[2])
        
        # Special handling for numbered lists - force white color for Google Slides compatibility
        # This runs after all text runs are created and formatted
        if bullet_data.get('type') == 'numbered':
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)  # Force white for numbered lists

def main():
    # Define paths
    template_path = Path(__file__).parent.parent / "template" / "template.pptx"
    output_path = Path(__file__).parent.parent / "output"
    script_path = output_path / "slide_script.json"
    
    # Create output directory if it doesn't exist
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Load the master template
    print(f"Loading master template from: {template_path}")
    prs = Presentation(str(template_path))
    
    # Map slide types to layouts based on template slides
    # Slide 1 (index 0) = cover layout
    # Slide 2 (index 1) = agenda layout  
    # Slide 3 (index 2) = subtitle/section_divider layout
    # Slide 4 (index 3) = content layout
    
    # Get the layouts from existing slides
    layout_map = {}
    if len(prs.slides) >= 4:
        layout_map['cover'] = prs.slides[0].slide_layout
        layout_map['agenda_slide'] = prs.slides[1].slide_layout
        layout_map['subtitle'] = prs.slides[2].slide_layout
        layout_map['section_divider'] = prs.slides[2].slide_layout  # Same as subtitle
        layout_map['content'] = prs.slides[3].slide_layout
        print(f"Identified layouts from template slides:")
        for slide_type, layout in layout_map.items():
            print(f"  {slide_type}: {layout.name}")
    else:
        print(f"Warning: Template has only {len(prs.slides)} slides, expected 4")
        # Fallback to using available layouts
        for i, slide in enumerate(prs.slides):
            print(f"  Slide {i+1}: {slide.slide_layout.name}")
    
    # Store original template slides count to delete them later
    original_slide_count = len(prs.slides)
    print(f"\nTemplate has {original_slide_count} original slides that will be removed later")
    
    # Load slide script
    print(f"\nLoading script from: {script_path}")
    with open(script_path, 'r', encoding='utf-8') as f:
        slides_data = json.load(f)
    print(f"Loaded {len(slides_data)} slides from script")
    
    # Create new slides with appropriate layouts and content
    print("\nCreating new slides with content...")
    for i, slide_data in enumerate(slides_data, 1):
        slide_type = slide_data.get('slide_type', 'content')
        print(f"  Slide {i}/{len(slides_data)}: {slide_type}")
        
        # Get the appropriate layout
        if slide_type in layout_map:
            layout = layout_map[slide_type]
        else:
            # Default to content layout if type not found
            print(f"    Warning: Unknown slide type '{slide_type}', using content layout")
            layout = layout_map.get('content', prs.slide_layouts[-1])
        
        # Add new slide with the correct layout
        new_slide = prs.slides.add_slide(layout)
        
        # Copy all content from the corresponding template slide
        # Find which template slide to copy from based on slide_type
        template_slide_idx = None
        if slide_type == 'cover':
            template_slide_idx = 0
        elif slide_type == 'agenda_slide':
            template_slide_idx = 1
        elif slide_type in ['subtitle', 'section_divider']:
            template_slide_idx = 2
        elif slide_type == 'content':
            template_slide_idx = 3
        
        # Copy shapes from template slide if we have it
        if template_slide_idx is not None and template_slide_idx < original_slide_count:
            template_slide = prs.slides[template_slide_idx]
            
            # Copy all non-placeholder shapes (images, graphics, etc.)
            for shape in template_slide.shapes:
                # Skip placeholders and empty text boxes
                try:
                    if shape.placeholder_format:
                        continue
                except (AttributeError, ValueError):
                    # Not a placeholder, continue with copying
                    pass
                    
                # Copy pictures/images
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        # Get image properties
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        # Get image bytes
                        image = shape.image
                        image_bytes = image.blob
                        
                        # Add picture to new slide
                        new_slide.shapes.add_picture(
                            io.BytesIO(image_bytes), left, top, width, height
                        )
                    except Exception as e:
                        print(f"    Could not copy image: {e}")
                
                # Copy other shapes (rectangles, lines, etc.) but not text
                elif not hasattr(shape, 'text_frame'):
                    # This would be more complex - shapes like rectangles, arrows, etc.
                    # For now, we'll skip these as they require more complex copying
                    pass
        
        # Remove only empty placeholder text shapes from the new slide
        shapes_to_remove = []
        for shape in new_slide.shapes:
            # Only remove text placeholders that are empty
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                try:
                    if shape.placeholder_format:
                        if not shape.text or shape.text.strip() == '':
                            shapes_to_remove.append(shape)
                except (AttributeError, ValueError):
                    # Not a placeholder, check if it's an empty text box
                    if not shape.text or shape.text.strip() == '':
                        shapes_to_remove.append(shape)
        
        # Remove the collected empty placeholders
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
        
        # Add content from JSON
        shapes_data = slide_data.get('shapes', [])
        for shape_data in shapes_data:
            if shape_data.get('type') == 'text':
                add_text_to_slide(new_slide, shape_data)
    
    # Delete the original template slides
    print(f"\nRemoving {original_slide_count} original template slides...")
    
    # We need to delete from the beginning, but deletion shifts indices
    # So we delete the first slide 'original_slide_count' times
    for i in range(original_slide_count):
        # Get the XML slide ID list
        xml_slides = prs.slides._sldIdLst
        # Remove the first slide (always index 0 after each deletion)
        slides = list(xml_slides)
        if slides:
            xml_slides.remove(slides[0])
            print(f"  Removed template slide {i+1}/{original_slide_count}")
    
    # Save the final presentation
    output_file = output_path / "final_presentation.pptx"
    print(f"\nSaving final presentation to: {output_file}")
    prs.save(str(output_file))
    
    print(f"\nDone! Created presentation with {len(prs.slides)} slides")
    print(f"All slides use proper layouts from the master template with backgrounds preserved.")

if __name__ == "__main__":
    main()