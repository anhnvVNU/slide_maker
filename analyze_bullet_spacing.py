#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu

def analyze_bullet_spacing_details(pptx_path, slide_index=3):
    """Phân tích chi tiết spacing giữa các bullet points"""
    print(f"Analyzing bullet spacing in slide {slide_index + 1} from: {pptx_path}")
    
    # Load presentation
    prs = Presentation(str(pptx_path))
    
    if slide_index >= len(prs.slides):
        print(f"Error: Slide {slide_index + 1} not found.")
        return
    
    slide = prs.slides[slide_index]
    
    print(f"\nSlide {slide_index + 1} Layout: {slide.slide_layout.name}")
    print("=" * 80)
    
    # Find main content shape (position x=626250)
    main_content_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'left') and int(shape.left) == 626250:
            main_content_shape = shape
            break
    
    if not main_content_shape or not hasattr(main_content_shape, 'text_frame'):
        print("Main content shape (x=626250) not found or has no text frame")
        return
    
    text_frame = main_content_shape.text_frame
    paragraphs = text_frame.paragraphs
    
    print(f"Main Content Shape Analysis:")
    print(f"Position: x={int(main_content_shape.left)}, y={int(main_content_shape.top)}")
    print(f"Size: {int(main_content_shape.width)} x {int(main_content_shape.height)}")
    print(f"Paragraphs count: {len(paragraphs)}")
    print()
    
    # Analyze each paragraph's spacing properties
    analysis = {
        "file": str(pptx_path),
        "slide": slide_index + 1,
        "total_paragraphs": len(paragraphs),
        "paragraphs": []
    }
    
    for i, para in enumerate(paragraphs):
        para_info = {
            "index": i,
            "text": para.text,
            "level": para.level,
            "char_count": len(para.text)
        }
        
        # Line spacing analysis
        try:
            if para.line_spacing:
                if hasattr(para.line_spacing, 'pt'):
                    para_info["line_spacing"] = {
                        "type": "points",
                        "value": float(para.line_spacing.pt)
                    }
                else:
                    para_info["line_spacing"] = {
                        "type": "multiple", 
                        "value": float(para.line_spacing)
                    }
            else:
                para_info["line_spacing"] = None
        except:
            para_info["line_spacing"] = "error"
        
        # Space before analysis
        try:
            if para.space_before:
                para_info["space_before"] = {
                    "pt": float(para.space_before.pt),
                    "emu": int(para.space_before)
                }
            else:
                para_info["space_before"] = None
        except:
            para_info["space_before"] = "error"
        
        # Space after analysis  
        try:
            if para.space_after:
                para_info["space_after"] = {
                    "pt": float(para.space_after.pt),
                    "emu": int(para.space_after)
                }
            else:
                para_info["space_after"] = None
        except:
            para_info["space_after"] = "error"
        
        # Bullet analysis via XML
        try:
            from lxml import etree
            p_elem = para._p
            pPr = p_elem.find('.//a:pPr', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            
            if pPr is not None:
                para_info["xml_properties"] = {}
                
                # Get all spacing-related attributes
                indent = pPr.get('indent')
                marL = pPr.get('marL')
                marR = pPr.get('marR') 
                marT = pPr.get('marT')
                marB = pPr.get('marB')
                
                if indent:
                    para_info["xml_properties"]["indent"] = {
                        "emu": int(indent),
                        "pt": int(indent) / 12700,
                        "inches": int(indent) / 914400
                    }
                
                if marL:
                    para_info["xml_properties"]["marL"] = {
                        "emu": int(marL),
                        "pt": int(marL) / 12700,
                        "inches": int(marL) / 914400
                    }
                
                if marR:
                    para_info["xml_properties"]["marR"] = {
                        "emu": int(marR),
                        "pt": int(marR) / 12700
                    }
                    
                if marT:
                    para_info["xml_properties"]["marT"] = {
                        "emu": int(marT),
                        "pt": int(marT) / 12700
                    }
                    
                if marB:
                    para_info["xml_properties"]["marB"] = {
                        "emu": int(marB),  
                        "pt": int(marB) / 12700
                    }
                
                # Check bullet character
                buChar = pPr.find('.//a:buChar', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if buChar is not None:
                    para_info["xml_properties"]["bullet_char"] = buChar.get('char')
        
        except Exception as e:
            para_info["xml_error"] = str(e)
        
        # Font size analysis
        if para.runs:
            run = para.runs[0]
            if run.font.size:
                para_info["font_size_pt"] = float(run.font.size.pt)
        
        analysis["paragraphs"].append(para_info)
        
        # Print detailed info
        bullet_char = ""
        if para_info.get("xml_properties", {}).get("bullet_char"):
            bullet_char = f" [{para_info['xml_properties']['bullet_char']}]"
        
        print(f"Para {i}: Level {para.level}{bullet_char}")
        print(f"  Text: '{para.text[:60]}{'...' if len(para.text) > 60 else ''}'")
        
        if para_info.get("font_size_pt"):
            print(f"  Font size: {para_info['font_size_pt']}pt")
        
        if para_info.get("line_spacing"):
            ls = para_info["line_spacing"]
            print(f"  Line spacing: {ls['value']} ({ls['type']})")
        
        if para_info.get("space_before"):
            sb = para_info["space_before"]
            print(f"  Space before: {sb['pt']}pt")
        
        if para_info.get("space_after"):
            sa = para_info["space_after"]
            print(f"  Space after: {sa['pt']}pt")
        
        if para_info.get("xml_properties"):
            xml = para_info["xml_properties"]
            spacing_attrs = []
            if "indent" in xml:
                spacing_attrs.append(f"indent={xml['indent']['pt']:.1f}pt")
            if "marL" in xml:
                spacing_attrs.append(f"marL={xml['marL']['pt']:.1f}pt")
            if "marT" in xml:
                spacing_attrs.append(f"marT={xml['marT']['pt']:.1f}pt") 
            if "marB" in xml:
                spacing_attrs.append(f"marB={xml['marB']['pt']:.1f}pt")
            
            if spacing_attrs:
                print(f"  XML spacing: {', '.join(spacing_attrs)}")
        
        print()
    
    # Calculate spacing between bullets
    print("SPACING ANALYSIS BETWEEN BULLETS:")
    print("-" * 50)
    
    effective_spacing = []
    for i in range(len(analysis["paragraphs"]) - 1):
        current = analysis["paragraphs"][i]
        next_para = analysis["paragraphs"][i + 1]
        
        # Calculate effective spacing between this bullet and next
        spacing_components = []
        total_spacing_pt = 0
        
        # Current paragraph's space_after
        if current.get("space_after"):
            space_after = current["space_after"]["pt"]
            spacing_components.append(f"space_after={space_after}pt")
            total_spacing_pt += space_after
        
        # Next paragraph's space_before
        if next_para.get("space_before"):
            space_before = next_para["space_before"]["pt"]
            spacing_components.append(f"space_before={space_before}pt")
            total_spacing_pt += space_before
        
        # Line spacing contribution
        if current.get("line_spacing"):
            ls = current["line_spacing"]
            if ls["type"] == "multiple":
                # Multiple line spacing creates internal spacing within paragraph
                line_height = (current.get("font_size_pt", 14) * ls["value"])
                spacing_components.append(f"line_height={line_height:.1f}pt")
            elif ls["type"] == "points":
                spacing_components.append(f"line_spacing={ls['value']}pt")
        
        effective_spacing.append({
            "between": f"{i} -> {i+1}",
            "components": spacing_components,
            "total_pt": total_spacing_pt
        })
        
        print(f"Between bullet {i} and {i+1}:")
        print(f"  Components: {', '.join(spacing_components) if spacing_components else 'None'}")
        print(f"  Total spacing: {total_spacing_pt}pt")
        print()
    
    analysis["effective_spacing"] = effective_spacing
    
    return analysis

def compare_bullet_spacing(template_path, output_path, slide_index=3):
    """So sánh spacing giữa template và output"""
    print("COMPARING BULLET SPACING:")
    print("=" * 80)
    
    template_analysis = analyze_bullet_spacing_details(template_path, slide_index)
    print("\n" + "="*80)
    output_analysis = analyze_bullet_spacing_details(output_path, slide_index)
    
    if not template_analysis or not output_analysis:
        print("Error: Could not analyze one or both files")
        return
    
    print("\n" + "="*80)
    print("SPACING COMPARISON SUMMARY:")
    print("="*80)
    
    template_spacing = template_analysis.get("effective_spacing", [])
    output_spacing = output_analysis.get("effective_spacing", [])
    
    min_comparisons = min(len(template_spacing), len(output_spacing))
    
    for i in range(min_comparisons):
        t_spacing = template_spacing[i]
        o_spacing = output_spacing[i]
        
        t_total = t_spacing["total_pt"]
        o_total = o_spacing["total_pt"]
        
        difference = o_total - t_total
        status = "✅ Same" if abs(difference) < 0.1 else f"❌ Diff: {difference:+.1f}pt"
        
        print(f"\nBetween bullets {i} and {i+1}:")
        print(f"  Template: {t_total:.1f}pt ({', '.join(t_spacing['components'])})")
        print(f"  Output:   {o_total:.1f}pt ({', '.join(o_spacing['components'])})")
        print(f"  Result:   {status}")
    
    # Save detailed analysis
    with open("template_bullet_spacing.json", 'w', encoding='utf-8') as f:
        json.dump(template_analysis, f, ensure_ascii=False, indent=2)
    
    with open("output_bullet_spacing.json", 'w', encoding='utf-8') as f:
        json.dump(output_analysis, f, ensure_ascii=False, indent=2)
    
    print(f"\nDetailed analysis saved to template_bullet_spacing.json and output_bullet_spacing.json")

def main():
    template_path = Path("/media/vunv/DATA/slide_maker_4/slide_maker_2/VuNV/slide_maker/template/20250813.pptx")
    output_path = Path("/media/vunv/DATA/slide_maker_4/slide_maker_2/VuNV/slide_maker/output/final_presentation.pptx")
    
    if not template_path.exists():
        print(f"Error: Template file not found: {template_path}")
        return
    
    if not output_path.exists():
        print(f"Error: Output file not found: {output_path}")
        return
    
    compare_bullet_spacing(template_path, output_path, slide_index=3)

if __name__ == "__main__":
    main()

